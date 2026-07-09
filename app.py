import os
import zipfile
import pythoncom
from flask import Flask, render_template, request, send_file

# ==============================================================================
# IMPORTAÇÃO SEGURA DE BIBLIOTECAS
# ==============================================================================
try:
    from PIL import Image
except ImportError:
    Image = None

try:
    from pypdf import PdfReader, PdfWriter
except ImportError:
    PdfReader = None
    PdfWriter = None

try:
    from moviepy import VideoFileClip
except ImportError:
    try:
        from moviepy.editor import VideoFileClip
    except ImportError:
        VideoFileClip = None

try:
    from pdf2docx import Converter
except ImportError:
    Converter = None

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    from openpyxl import Workbook
except ImportError:
    Workbook = None

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    Presentation = None
    Inches = None

try:
    from docx2pdf import convert as docx_to_pdf_convert
except ImportError:
    docx_to_pdf_convert = None

try:
    import win32com.client
except ImportError:
    win32com = None

# ==============================================================================
# CONFIGURAÇÃO DO APP
# ==============================================================================
app = Flask(__name__)
UPLOAD_FOLDER = 'uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Interface Base unificada com o novo visual "Converte Ai" (Clean, Azul e Monetizável)
def layout_base(conteudo_pagina, ferramenta_ativa):
    menu_items = {
        'imagem-to-pdf': ('/', 'fa-file-image', 'Imagem para PDF'),
        'juntar': ('/juntar-view', 'fa-file-pdf', 'Juntar PDFs'),
        'dividir': ('/dividir-view', 'fa-scissors', 'Dividir PDF'),
        'pdf-to-word': ('/pdf-to-word-view', 'fa-file-word', 'PDF para Word'),
        'pdf-to-excel': ('/pdf-to-excel-view', 'fa-file-excel', 'PDF para Excel'),
        'pdf-to-pptx': ('/pdf-to-pptx-view', 'fa-file-powerpoint', 'PDF para PPTX'),
        'pdf-to-image': ('/pdf-to-image-view', 'fa-images', 'PDF para Imagem'),
        'word-to-pdf': ('/word-to-pdf-view', 'fa-file-word', 'Word para PDF'),
        'excel-to-pdf': ('/excel-to-pdf-view', 'fa-file-excel', 'Excel para PDF'),
        'pptx-to-pdf': ('/pptx-to-pdf-view', 'fa-file-powerpoint', 'PPTX para PDF'),
        'mp4-to-mp3': ('/mp4-to-mp3-view', 'fa-music', 'MP4 para MP3'),
    }
    
    categorias = {
        'Modificar PDF': ['juntar', 'dividir'],
        'Converter de PDF': ['pdf-to-word', 'pdf-to-excel', 'pdf-to-pptx', 'pdf-to-image'],
        'Converter para PDF': ['word-to-pdf', 'excel-to-pdf', 'pptx-to-pdf', 'imagem-to-pdf'],
        'Áudio & Vídeo': ['mp4-to-mp3']
    }
    
    sidebar_html = ""
    for cat_name, keys in categorias.items():
        sidebar_html += f'<div class="menu-category">{cat_name}</div>'
        for key in keys:
            route, icon, label = menu_items[key]
            active_class = "active" if ferramenta_ativa == key else ""
            sidebar_html += f'<a href="{route}" class="menu-item {active_class}"><i class="fa-solid {icon}"></i>{label}</a>'
    
    return f'''
    <!DOCTYPE html>
    <html lang="pt-br">
    <head>
        <meta charset="UTF-8">
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <title>Converte Ai - Conversor de Arquivos Online Grátis</title>
        <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.4.0/css/all.min.css">
        <style>
            :root {{
                --primary-dark: #0f172a; /* Azul Escuro / Slate */
                --primary: #1d4ed8;      /* Azul Principal */
                --primary-light: #eff6ff;/* Azul Claro de Fundo */
                --accent: #2563eb;       /* Azul de Destaque */
                --bg-main: #f8fafc;
                --bg-card: #ffffff;
                --text-dark: #1e293b;
                --text-muted: #64748b;
                --border-color: #e2e8f0;
                --sidebar-width: 260px;
                --ad-width: 180px;
            }}
            body {{
                font-family: 'Segoe UI', system-ui, -apple-system, sans-serif;
                margin: 0; padding: 0;
                background-color: var(--bg-main);
                color: var(--text-dark);
                display: flex;
                min-height: 100vh;
            }}
            /* BARRA VERTICAL LATERAL */
            .sidebar {{
                width: var(--sidebar-width);
                background-color: var(--primary-dark);
                color: #f1f5f9;
                display: flex; flex-direction: column;
                position: fixed;
                height: 100vh;
                overflow-y: auto;
                z-index: 100;
            }}
            .sidebar-header {{
                padding: 24px;
                display: flex; align-items: center; gap: 10px;
                border-bottom: 1px solid #1e293b;
            }}
            .sidebar-header i {{ font-size: 24px; color: #60a5fa; }}
            .sidebar-header h1 {{ font-size: 20px; margin: 0; font-weight: 800; color: #ffffff; letter-spacing: -0.5px; }}
            .sidebar-header h1 span {{ color: #3b82f6; }}
            .menu-category {{
                font-size: 11px; text-transform: uppercase; letter-spacing: 0.5px;
                color: #475569; padding: 16px 24px 6px; font-weight: 700;
            }}
            .menu-item {{
                display: flex; align-items: center; gap: 12px;
                padding: 10px 24px; color: #94a3b8; text-decoration: none;
                font-size: 14px; font-weight: 500; transition: all 0.15s ease;
            }}
            .menu-item:hover {{ background-color: #1e293b; color: #f8fafc; }}
            .menu-item.active {{
                background-color: var(--primary);
                color: #ffffff; border-radius: 6px;
                margin: 2px 12px; padding: 10px 12px;
            }}
            .menu-item i {{ font-size: 15px; width: 20px; text-align: center; }}
            
            /* ÁREA PRINCIPAL */
            .main-wrapper {{
                margin-left: var(--sidebar-width);
                margin-right: var(--ad-width); /* Espaço para o anúncio da direita */
                flex-grow: 1;
                display: flex; flex-direction: column;
                min-height: 100vh;
                box-sizing: border-box;
            }}
            
            /* ÁREA DE APRESENTAÇÃO (HERO SECTION) */
            .hero-section {{
                background-color: #ffffff;
                border-bottom: 1px solid var(--border-color);
                padding: 40px; text-align: left;
            }}
            .hero-container {{ max-width: 800px; margin: 0 auto; }}
            .hero-section h2 {{ margin: 0 0 10px 0; font-size: 28px; font-weight: 800; color: var(--primary-dark); }}
            .hero-section p {{ margin: 0; font-size: 15px; color: var(--text-muted); line-height: 1.6; }}
            
            /* CONTEÚDO DA FERRAMENTA */
            .content-area {{
                flex-grow: 1;
                display: flex; flex-direction: column; align-items: center; justify-content: center;
                padding: 40px; box-sizing: border-box;
            }}
            .tool-container {{
                background: var(--bg-card);
                padding: 40px; border-radius: 12px;
                box-shadow: 0 1px 3px rgba(0,0,0,0.05);
                border: 1px solid var(--border-color);
                max-width: 460px; width: 100%; text-align: center;
            }}
            .tool-title {{ color: var(--primary-dark); margin: 0 0 8px 0; font-size: 20px; font-weight: 700; }}
            .tool-desc {{ color: var(--text-muted); font-size: 13px; margin-bottom: 24px; line-height: 1.4; }}
            
            /* DROPZONE CLEAN */
            .file-dropzone {{
                border: 2px dashed #cbd5e1; background-color: var(--primary-light);
                border-radius: 10px; padding: 30px 20px; margin-bottom: 24px;
                cursor: pointer; transition: all 0.2s ease; position: relative;
            }}
            .file-dropzone:hover {{ border-color: var(--primary); background-color: #e0f2fe; }}
            .file-dropzone i {{ font-size: 36px; color: var(--primary); margin-bottom: 10px; }}
            .file-dropzone p {{ margin: 0; font-size: 14px; color: #334155; font-weight: 500; }}
            .file-dropzone input[type="file"] {{ position: absolute; top: 0; left: 0; width: 100%; height: 100%; opacity: 0; cursor: pointer; }}
            .file-name-display {{ margin-top: 8px; font-size: 13px; color: var(--primary); font-weight: 600; word-break: break-all; }}
            
            /* BOTÃO AZUL */
            button {{
                background-color: var(--primary); color: white; border: none; padding: 12px 24px;
                border-radius: 6px; cursor: pointer; font-size: 14px; font-weight: 600; width: 100%;
                transition: background 0.15s ease;
            }}
            button:hover {{ background-color: var(--accent); }}
            
            /* ESPAÇOS ESTRUTURADOS PARA ANÚNCIOS (ADS) */
            .ad-sidebar-right {{
                position: fixed; right: 0; top: 0; width: var(--ad-width); height: 100vh;
                background-color: #f1f5f9; border-left: 1px solid var(--border-color);
                display: flex; align-items: center; justify-content: center; z-index: 90;
            }}
            .ad-banner-bottom {{
                width: 100%; padding: 20px 0; background-color: #f1f5f9;
                border-top: 1px solid var(--border-color); text-align: center; margin-top: auto;
            }}
            .ad-box {{
                background: #e2e8f0; border: 1px dashed #cbd5e1; color: #94a3b8;
                font-size: 11px; font-weight: 600; text-transform: uppercase; letter-spacing: 1px;
                display: flex; align-items: center; justify-content: center; border-radius: 4px;
            }}
            .ad-vertical {{ width: 160px; height: 600px; }}
            .ad-horizontal {{ width: 728px; height: 90px; margin: 0 auto; }}
            
            /* LOADER */
            .loader-container {{ display: none; margin-top: 20px; }}
            .loader {{ border: 3px solid #f3f3f3; border-top: 3px solid var(--primary); border-radius: 50%; width: 24px; height: 24px; animation: spin 1s linear infinite; margin: 0 auto 8px; }}
            @keyframes spin {{ 0% {{ transform: rotate(0deg); }} 100% {{ transform: rotate(360deg); }} }}
            .loader-text {{ font-size: 12px; color: var(--text-muted); }}
        </style>
        <script>
            function handleFileSelect(input, displayId) {{
                const display = document.getElementById(displayId);
                if (input.files && input.files.length > 0) {{
                    display.textContent = input.files.length === 1 ? "Selecionado: " + input.files[0].name : input.files.length + " arquivos selecionados";
                }} else {{ display.textContent = ""; }}
            }}
            function showLoader() {{
                document.getElementById('loader').style.display = 'block';
                document.getElementById('submit-btn').style.display = 'none';
            }}
        </script>
    </head>
    <body>
        <!-- BARRA LATERAL ESQUERDA -->
        <div class="sidebar">
            <div class="sidebar-header">
                <i class="fa-solid fa-bolt"></i>
                <h1>Converte<span>Ai</span></h1>
            </div>
            <div class="sidebar-menu">
                {sidebar_html}
            </div>
        </div>
        
        <!-- ÁREA CENTRAL DE CONTEÚDO -->
        <div class="main-wrapper">
            <!-- ÁREA DE APRESENTAÇÃO (HERO) -->
            <div class="hero-section">
                <div class="hero-container">
                    <h2>Conversão Inteligente de Arquivos</h2>
                    <p>Bem-vindo ao Converte Ai. Uma plataforma web limpa, rápida e totalmente gratuita projetada para simplificar sua rotina. Converta, junte ou divida seus documentos e arquivos multimídia diretamente no navegador, de forma 100% segura e sem instalações.</p>
                </div>
            </div>
            
            <!-- ÁREA DA FERRAMENTA ATIVA -->
            <div class="content-area">
                <div class="tool-container">
                    {conteudo_pagina}
                    <div id="loader" class="loader-container">
                        <div class="loader"></div>
                        <div class="loader-text">Processando arquivo...</div>
                    </div>
                </div>
            </div>
            
            <!-- ESPAÇO PARA ANÚNCIO INFERIOR (BANNER HORIZONTAL) -->
            <div class="ad-banner-bottom">
                <div class="ad-box ad-horizontal">Espaço para Anúncio (728x90)</div>
            </div>
        </div>
        
        <!-- ESPAÇO PARA ANÚNCIO LATERAL DIREITO (ARRANHA-CÉU) -->
        <div class="ad-sidebar-right">
            <div class="ad-box ad-vertical">Espaço para Anúncio (160x600)</div>
        </div>
    </body>
    </html>
    '''

# ==============================================================================
# ROTA 1: IMAGEM PARA PDF (HOME)
# ==============================================================================
@app.route('/')
def home():
    conteudo = '''
        <h3 class="tool-title">Imagem para PDF</h3>
        <p class="tool-desc">Converta suas imagens JPG ou PNG em um documento PDF instantâneo.</p>
        <form action="/convert" method="post" enctype="multipart/form-data" onsubmit="showLoader()">
            <div class="file-dropzone">
                <i class="fa-solid fa-file-image"></i>
                <p>Arraste a imagem ou clique para selecionar</p>
                <input type="file" name="imagem_usuario" accept=".jpg, .jpeg, .png" required onchange="handleFileSelect(this, 'file-name')">
                <div id="file-name" class="file-name-display"></div>
            </div>
            <button type="submit" id="submit-btn">Converter para PDF</button>
        </form>
    '''
    return layout_base(conteudo, 'imagem-to-pdf')

@app.route('/convert', methods=['POST'])
def convert_file():
    if Image is None: return "Erro: Biblioteca 'Pillow' não instalada.", 500
    if 'imagem_usuario' not in request.files: return 'Nenhuma imagem enviada', 400
    arquivo = request.files['imagem_usuario']
    if arquivo.filename == '': return 'Nome inválido', 400
    
    caminho_imagem = os.path.join(UPLOAD_FOLDER, arquivo.filename)
    arquivo.save(caminho_imagem)
    
    try:
        img = Image.open(caminho_imagem)
        img_rgb = img.convert('RGB')
        nome_pdf = os.path.splitext(arquivo.filename)[0] + '.pdf'
        caminho_pdf = os.path.join(UPLOAD_FOLDER, nome_pdf)
        img_rgb.save(caminho_pdf, 'PDF', resolution=100.0)
        return send_file(caminho_pdf, as_attachment=True, download_name=nome_pdf)
    except Exception as e:
        return f'Erro: {str(e)}', 500
    finally:
        if os.path.exists(caminho_imagem): os.remove(caminho_imagem)

# ==============================================================================
# ROTA 2: JUNTAR PDFs
# ==============================================================================
@app.route('/juntar-view')
def juntar_view():
    conteudo = '''
        <h3 class="tool-title">Juntar PDFs</h3>
        <p class="tool-desc">Combine dois ou mais arquivos em um único documento estruturado.</p>
        <form action="/juntar" method="post" enctype="multipart/form-data" onsubmit="showLoader()">
            <div class="file-dropzone">
                <i class="fa-solid fa-file-pdf"></i>
                <p>Selecione os múltiplos arquivos PDF</p>
                <input type="file" name="arquivos_pdf" accept=".pdf" multiple required onchange="handleFileSelect(this, 'file-name')">
                <div id="file-name" class="file-name-display"></div>
            </div>
            <button type="submit" id="submit-btn">Juntar PDFs</button>
        </form>
    '''
    return layout_base(conteudo, 'juntar')

@app.route('/juntar', methods=['POST'])
def juntar_pdfs():
    if PdfWriter is None: return "Erro: Biblioteca 'pypdf' não instalada.", 500
    arquivos = request.files.getlist('arquivos_pdf')
    if not arquivos or arquivos[0].filename == '': return 'Nenhum arquivo selecionado', 400
    
    merger = PdfWriter()
    caminhos_temporarios = []
    try:
        for arquivo in arquivos:
            caminho_temp = os.path.join(UPLOAD_FOLDER, arquivo.filename)
            arquivo.save(caminho_temp)
            caminhos_temporarios.append(caminho_temp)
            merger.append(caminho_temp)
        
        nome_saida = "pdfs_combinados.pdf"
        caminho_saida = os.path.join(UPLOAD_FOLDER, nome_saida)
        with open(caminho_saida, "wb") as f:
            merger.write(f)
        merger.close()
        return send_file(caminho_saida, as_attachment=True, download_name=nome_saida)
    except Exception as e:
        return f'Erro: {str(e)}', 500
    finally:
        for c in caminhos_temporarios:
            if os.path.exists(c): os.remove(c)

# ==============================================================================
# ROTA 3: DIVIDIR PDF
# ==============================================================================
@app.route('/dividir-view')
def dividir_view():
    conteudo = '''
        <h3 class="tool-title">Dividir PDF</h3>
        <p class="tool-desc">Separe todas as páginas do seu documento PDF em arquivos individuais dentro de um ZIP.</p>
        <form action="/dividir" method="post" enctype="multipart/form-data" onsubmit="showLoader()">
            <div class="file-dropzone">
                <i class="fa-solid fa-scissors"></i>
                <p>Escolha o arquivo PDF para fatiar</p>
                <input type="file" name="arquivo_pdf" accept=".pdf" required onchange="handleFileSelect(this, 'file-name')">
                <div id="file-name" class="file-name-display"></div>
            </div>
            <button type="submit" id="submit-btn">Dividir PDF</button>
        </form>
    '''
    return layout_base(conteudo, 'dividir')

@app.route('/dividir', methods=['POST'])
def dividir_pdf():
    if PdfReader is None: return "Erro: Biblioteca 'pypdf' não instalada.", 500
    arquivo = request.files['arquivo_pdf']
    if arquivo.filename == '': return 'Arquivo inválido', 400
    
    caminho_pdf = os.path.join(UPLOAD_FOLDER, arquivo.filename)
    arquivo.save(caminho_pdf)
    
    try:
        reader = PdfReader(caminho_pdf)
        nome_base = os.path.splitext(arquivo.filename)[0]
        nome_zip = f"{nome_base}_dividido.zip"
        caminho_zip = os.path.join(UPLOAD_FOLDER, nome_zip)
        
        with zipfile.ZipFile(caminho_zip, 'w') as z:
            for idx, pagina in enumerate(reader.pages):
                writer = PdfWriter()
                writer.add_page(pagina)
                nome_pagina = f"{nome_base}_pag_{idx+1}.pdf"
                caminho_pagina = os.path.join(UPLOAD_FOLDER, nome_pagina)
                with open(caminho_pagina, "wb") as out:
                    writer.write(out)
                z.write(caminho_pagina, nome_pagina)
                os.remove(caminho_pagina)
                
        return send_file(caminho_zip, as_attachment=True, download_name=nome_zip)
    except Exception as e:
        return f'Erro: {str(e)}', 500
    finally:
        if os.path.exists(caminho_pdf): os.remove(caminho_pdf)

# ==============================================================================
# ROTA 4: PDF PARA WORD
# ==============================================================================
@app.route('/pdf-to-word-view')
def pdf_to_word_view():
    conteudo = '''
        <h3 class="tool-title">PDF para Word</h3>
        <p class="tool-desc">Converta o conteúdo do seu PDF para um documento editável do Word (.docx).</p>
        <form action="/pdf-to-word" method="post" enctype="multipart/form-data" onsubmit="showLoader()">
            <div class="file-dropzone">
                <i class="fa-solid fa-file-word"></i>
                <p>Selecione o arquivo PDF</p>
                <input type="file" name="arquivo_pdf" accept=".pdf" required onchange="handleFileSelect(this, 'file-name')">
                <div id="file-name" class="file-name-display"></div>
            </div>
            <button type="submit" id="submit-btn">Converter para Word</button>
        </form>
    '''
    return layout_base(conteudo, 'pdf-to-word')

@app.route('/pdf-to-word', methods=['POST'])
def pdf_to_word():
    if Converter is None: return "Erro: Biblioteca 'pdf2docx' não instalada.", 500
    arquivo = request.files['arquivo_pdf']
    if arquivo.filename == '': return 'Arquivo inválido', 400
    
    caminho_pdf = os.path.join(UPLOAD_FOLDER, arquivo.filename)
    arquivo.save(caminho_pdf)
    nome_docx = os.path.splitext(arquivo.filename)[0] + '.docx'
    caminho_docx = os.path.join(UPLOAD_FOLDER, nome_docx)
    
    try:
        cv = Converter(caminho_pdf)
        cv.convert(caminho_docx, start=0, end=None)
        cv.close()
        return send_file(caminho_docx, as_attachment=True, download_name=nome_docx)
    except Exception as e:
        return f'Erro na conversão: {str(e)}', 500
    finally:
        if os.path.exists(caminho_pdf): os.remove(caminho_pdf)

# ==============================================================================
# ROTA 5: PDF PARA EXCEL
# ==============================================================================
@app.route('/pdf-to-excel-view')
def pdf_to_excel_view():
    conteudo = '''
        <h3 class="tool-title">PDF para Excel</h3>
        <p class="tool-desc">Extraia as matrizes de tabelas do PDF direto para abas organizadas em planilhas (.xlsx).</p>
        <form action="/pdf-to-excel" method="post" enctype="multipart/form-data" onsubmit="showLoader()">
            <div class="file-dropzone">
                <i class="fa-solid fa-file-excel"></i>
                <p>Selecione o arquivo PDF com tabelas</p>
                <input type="file" name="arquivo_pdf" accept=".pdf" required onchange="handleFileSelect(this, 'file-name')">
                <div id="file-name" class="file-name-display"></div>
            </div>
            <button type="submit" id="submit-btn">Converter para Excel</button>
        </form>
    '''
    return layout_base(conteudo, 'pdf-to-excel')

@app.route('/pdf-to-excel', methods=['POST'])
def pdf_to_excel():
    if pdfplumber is None or Workbook is None: return "Erro: Bibliotecas ausentes.", 500
    arquivo = request.files['arquivo_pdf']
    if arquivo.filename == '': return 'Arquivo inválido', 400
    
    caminho_pdf = os.path.join(UPLOAD_FOLDER, arquivo.filename)
    arquivo.save(caminho_pdf)
    nome_xlsx = os.path.splitext(arquivo.filename)[0] + '.xlsx'
    caminho_xlsx = os.path.join(UPLOAD_FOLDER, nome_xlsx)
    
    try:
        wb = Workbook()
        wb.remove(wb.active)
        
        with pdfplumber.open(caminho_pdf) as pdf:
            for idx, pagina in enumerate(pdf.pages):
                ws = wb.create_sheet(title=f"Pagina {idx+1}")
                tabelas = pagina.extract_tables()
                if tabelas:
                    for tabela in tabelas:
                        for linha in tabela:
                            linha_limpa = [str(c) if c is not None else "" for c in linha]
                            ws.append(linha_limpa)
                        ws.append([])
                else:
                    ws.append(["Nenhuma tabela estruturada encontrada nesta página."])
                    
        wb.save(caminho_xlsx)
        return send_file(caminho_xlsx, as_attachment=True, download_name=nome_xlsx)
    except Exception as e:
        return f'Erro: {str(e)}', 500
    finally:
        if os.path.exists(caminho_pdf): os.remove(caminho_pdf)

# ==============================================================================
# ROTA 6: PDF PARA PPTX
# ==============================================================================
@app.route('/pdf-to-pptx-view')
def pdf_to_pptx_view():
    conteudo = '''
        <h3 class="tool-title">PDF para PowerPoint</h3>
        <p class="tool-desc">Converta as páginas estruturadas de um PDF em slides prontos (.pptx).</p>
        <form action="/pdf-to-pptx" method="post" enctype="multipart/form-data" onsubmit="showLoader()">
            <div class="file-dropzone">
                <i class="fa-solid fa-file-powerpoint"></i>
                <p>Selecione o arquivo PDF original</p>
                <input type="file" name="arquivo_pdf" accept=".pdf" required onchange="handleFileSelect(this, 'file-name')">
                <div id="file-name" class="file-name-display"></div>
            </div>
            <button type="submit" id="submit-btn">Converter para PPTX</button>
        </form>
    '''
    return layout_base(conteudo, 'pdf-to-pptx')

@app.route('/pdf-to-pptx', methods=['POST'])
def pdf_to_pptx():
    if fitz is None or Presentation is None: return "Erro: Bibliotecas ausentes.", 500
    arquivo = request.files['arquivo_pdf']
    if arquivo.filename == '': return 'Arquivo inválido', 400
    
    caminho_pdf = os.path.join(UPLOAD_FOLDER, arquivo.filename)
    arquivo.save(caminho_pdf)
    nome_pptx = os.path.splitext(arquivo.filename)[0] + '.pptx'
    caminho_pptx = os.path.join(UPLOAD_FOLDER, nome_pptx)
    
    try:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        doc = fitz.open(caminho_pdf)
        for idx in range(len(doc)):
            pagina = doc.load_page(idx)
            pix = pagina.get_pixmap(dpi=150)
            
            caminho_temp_img = os.path.join(UPLOAD_FOLDER, f"temp_slide_{idx}.png")
            pix.save(caminho_temp_img)
            
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(caminho_temp_img, 0, 0, width=prs.slide_width, height=prs.slide_height)
            os.remove(caminho_temp_img)
            
        prs.save(caminho_pptx)
        return send_file(caminho_pptx, as_attachment=True, download_name=nome_pptx)
    except Exception as e:
        return f'Erro: {str(e)}', 500
    finally:
        if os.path.exists(caminho_pdf): os.remove(caminho_pdf)

# ==============================================================================
# ROTA 7: PDF PARA IMAGEM
# ==============================================================================
@app.route('/pdf-to-image-view')
def pdf_to_image_view():
    conteudo = '''
        <h3 class="tool-title">PDF para Imagem</h3>
        <p class="tool-desc">Extraia as folhas do seu arquivo PDF em formato de fotos avulsas de alta resolução (.png).</p>
        <form action="/pdf-to-image" method="post" enctype="multipart/form-data" onsubmit="showLoader()">
            <div class="file-dropzone">
                <i class="fa-solid fa-images"></i>
                <p>Selecione o arquivo PDF</p>
                <input type="file" name="arquivo_pdf" accept=".pdf" required onchange="handleFileSelect(this, 'file-name')">
                <div id="file-name" class="file-name-display"></div>
            </div>
            <button type="submit" id="submit-btn">Converter para Imagem</button>
        </form>
    '''
    return layout_base(conteudo, 'pdf-to-image')

@app.route('/pdf-to-image', methods=['POST'])
def pdf_to_image():
    if fitz is None: return "Erro: Biblioteca 'pymupdf' não instalada.", 500
    arquivo = request.files['arquivo_pdf']
    if arquivo.filename == '': return 'Arquivo inválido', 400
    
    caminho_pdf = os.path.join(UPLOAD_FOLDER, arquivo.filename)
    arquivo.save(caminho_pdf)
    nome_base = os.path.splitext(arquivo.filename)[0]
    
    try:
        doc = fitz.open(caminho_pdf)
        if len(doc) == 1:
            pix = doc[0].get_pixmap(dpi=150)
            nome_img = f"{nome_base}.png"
            caminho_img = os.path.join(UPLOAD_FOLDER, nome_img)
            pix.save(caminho_img)
            return send_file(caminho_img, as_attachment=True, download_name=nome_img)
        else:
            nome_zip = f"{nome_base}_imagens.zip"
            caminho_zip = os.path.join(UPLOAD_FOLDER, nome_zip)
            with zipfile.ZipFile(caminho_zip, 'w') as z:
                for idx, pagina in enumerate(doc):
                    pix = pagina.get_pixmap(dpi=150)
                    nome_img = f"{nome_base}_pag_{idx+1}.png"
                    caminho_img = os.path.join(UPLOAD_FOLDER, nome_img)
                    pix.save(caminho_img)
                    z.write(caminho_img, nome_img)
                    os.remove(caminho_img)
            return send_file(caminho_zip, as_attachment=True, download_name=nome_zip)
    except Exception as e:
        return f'Erro: {str(e)}', 500
    finally:
        if os.path.exists(caminho_pdf): os.remove(caminho_pdf)

# ==============================================================================
# ROTA 8: WORD PARA PDF
# ==============================================================================
@app.route('/word-to-pdf-view')
def word_to_pdf_view():
    conteudo = '''
        <h3 class="tool-title">Word para PDF</h3>
        <p class="tool-desc">Gere PDFs fiéis a partir de documentos de texto do Word (.docx).</p>
        <form action="/word-to-pdf" method="post" enctype="multipart/form-data" onsubmit="showLoader()">
            <div class="file-dropzone">
                <i class="fa-solid fa-file-pdf"></i>
                <p>Selecione o documento do Word (.docx)</p>
                <input type="file" name="arquivo_docx" accept=".docx" required onchange="handleFileSelect(this, 'file-name')">
                <div id="file-name" class="file-name-display"></div>
            </div>
            <button type="submit" id="submit-btn">Converter para PDF</button>
        </form>
    '''
    return layout_base(conteudo, 'word-to-pdf')

@app.route('/word-to-pdf', methods=['POST'])
def word_to_pdf():
    if docx_to_pdf_convert is None: return "Erro: Biblioteca 'docx2pdf' não instalada.", 500
    arquivo = request.files['arquivo_docx']
    if arquivo.filename == '': return 'Arquivo inválido', 400
    
    caminho_docx = os.path.join(UPLOAD_FOLDER, arquivo.filename)
    arquivo.save(caminho_docx)
    nome_pdf = os.path.splitext(arquivo.filename)[0] + '.pdf'
    caminho_pdf = os.path.join(UPLOAD_FOLDER, nome_pdf)
    
    try:
        pythoncom.CoInitialize()
        docx_to_pdf_convert(caminho_docx, caminho_pdf)
        return send_file(caminho_pdf, as_attachment=True, download_name=nome_pdf)
    except Exception as e:
        return f'Erro: {str(e)}. Verifique se possui o Microsoft Word instalado.', 500
    finally:
        if os.path.exists(caminho_docx): os.remove(caminho_docx)

# ==============================================================================
# ROTA 9: EXCEL PARA PDF
# ==============================================================================
@app.route('/excel-to-pdf-view')
def excel_to_pdf_view():
    conteudo = '''
        <h3 class="tool-title">Excel para PDF</h3>
        <p class="tool-desc">Transforme suas planilhas operacionais e tabelas do Excel (.xlsx) em formato PDF.</p>
        <form action="/excel-to-pdf" method="post" enctype="multipart/form-data" onsubmit="showLoader()">
            <div class="file-dropzone">
                <i class="fa-solid fa-file-pdf"></i>
                <p>Selecione a planilha Excel (.xlsx)</p>
                <input type="file" name="arquivo_xlsx" accept=".xlsx" required onchange="handleFileSelect(this, 'file-name')">
                <div id="file-name" class="file-name-display"></div>
            </div>
            <button type="submit" id="submit-btn">Converter para PDF</button>
        </form>
    '''
    return layout_base(conteudo, 'excel-to-pdf')

@app.route('/excel-to-pdf', methods=['POST'])
def excel_to_pdf():
    if win32com is None: return "Erro: pywin32 indisponível.", 500
    arquivo = request.files['arquivo_xlsx']
    if arquivo.filename == '': return 'Arquivo inválido', 400
    
    caminho_xlsx = os.path.join(UPLOAD_FOLDER, arquivo.filename)
    arquivo.save(caminho_xlsx)
    nome_pdf = os.path.splitext(arquivo.filename)[0] + '.pdf'
    caminho_pdf = os.path.join(UPLOAD_FOLDER, nome_pdf)
    
    try:
        pythoncom.CoInitialize()
        excel = win32com.client.Dispatch("Excel.Application")
        excel.Visible = False
        excel.DisplayAlerts = False
        
        wb = excel.Workbooks.Open(os.path.abspath(caminho_xlsx))
        wb.ExportAsFixedFormat(0, os.path.abspath(caminho_pdf))
        wb.Close(False)
        excel.Quit()
        return send_file(caminho_pdf, as_attachment=True, download_name=nome_pdf)
    except Exception as e:
        return f'Erro: {str(e)}. Verifique se possui o Microsoft Excel instalado.', 500
    finally:
        if os.path.exists(caminho_xlsx): os.remove(caminho_xlsx)

# ==============================================================================
# ROTA 10: PPTX PARA PDF
# ==============================================================================
@app.route('/pptx-to-pdf-view')
def pptx_to_pdf_view():
    conteudo = '''
        <h3 class="tool-title">PowerPoint para PDF</h3>
        <p class="tool-desc">Exporte seus slides de apresentações (.pptx) para um documento PDF estável.</p>
        <form action="/pptx-to-pdf" method="post" enctype="multipart/form-data" onsubmit="showLoader()">
            <div class="file-dropzone">
                <i class="fa-solid fa-file-pdf"></i>
                <p>Selecione a apresentação (.pptx)</p>
                <input type="file" name="arquivo_pptx" accept=".pptx" required onchange="handleFileSelect(this, 'file-name')">
                <div id="file-name" class="file-name-display"></div>
            </div>
            <button type="submit" id="submit-btn">Converter para PDF</button>
        </form>
    '''
    return layout_base(conteudo, 'pptx-to-pdf')

@app.route('/pptx-to-pdf', methods=['POST'])
def pptx_to_pdf():
    if win32com is None: return "Erro: pywin32 indisponível.", 500
    arquivo = request.files['arquivo_pptx']
    if arquivo.filename == '': return 'Arquivo inválido', 400
    
    caminho_pptx = os.path.join(UPLOAD_FOLDER, arquivo.filename)
    arquivo.save(caminho_pptx)
    nome_pdf = os.path.splitext(arquivo.filename)[0] + '.pdf'
    caminho_pdf = os.path.join(UPLOAD_FOLDER, nome_pdf)
    
    try:
        pythoncom.CoInitialize()
        powerpoint = win32com.client.Dispatch("Powerpoint.Application")
        
        ppt = powerpoint.Presentations.Open(os.path.abspath(caminho_pptx), WithWindow=False)
        ppt.SaveAs(os.path.abspath(caminho_pdf), 32)
        ppt.Close()
        powerpoint.Quit()
        return send_file(caminho_pdf, as_attachment=True, download_name=nome_pdf)
    except Exception as e:
        return f'Erro: {str(e)}. Verifique se possui o Microsoft PowerPoint instalado.', 500
    finally:
        if os.path.exists(caminho_pptx): os.remove(caminho_pptx)

# ==============================================================================
# ROTA 11: MP4 PARA MP3
# ==============================================================================
@app.route('/mp4-to-mp3-view')
def mp4_to_mp3_view():
    conteudo = '''
        <h3 class="tool-title">MP4 para MP3</h3>
        <p class="tool-desc">Extraia as faixas de áudio das suas mídias de vídeo e salve em MP3 leve.</p>
        <form action="/converter-audio" method="post" enctype="multipart/form-data" onsubmit="showLoader()">
            <div class="file-dropzone">
                <i class="fa-solid fa-music"></i>
                <p>Selecione o vídeo MP4</p>
                <input type="file" name="video_usuario" accept=".mp4" required onchange="handleFileSelect(this, 'file-name')">
                <div id="file-name" class="file-name-display"></div>
            </div>
            <button type="submit" id="submit-btn">Extrair Áudio (MP3)</button>
        </form>
    '''
    return layout_base(conteudo, 'mp4-to-mp3')

@app.route('/converter-audio', methods=['POST'])
def converter_audio():
    if VideoFileClip is None: return "Erro: Biblioteca MoviePy ausente.", 500
    arquivo = request.files['video_usuario']
    if arquivo.filename == '': return 'Arquivo inválido', 400
    
    caminho_video = os.path.join(UPLOAD_FOLDER, arquivo.filename)
    arquivo.save(caminho_video)
    nome_audio = os.path.splitext(arquivo.filename)[0] + '.mp3'
    caminho_audio = os.path.join(UPLOAD_FOLDER, nome_audio)
    
    try:
        video = VideoFileClip(caminho_video)
        video.audio.write_audiofile(caminho_audio, codec='mp3')
        video.close()
        return send_file(caminho_audio, as_attachment=True, download_name=nome_audio)
    except Exception as e:
        return f'Erro ao extrair áudio: {str(e)}', 500
    finally:
        if os.path.exists(caminho_video): os.remove(caminho_video)

if __name__ == '__main__':
    app.run(debug=True)